import os
import io
import base64
import mimetypes
import time
import requests
import validators
from urllib.parse import urlparse, unquote
from flask import Blueprint, jsonify, request
from werkzeug.utils import secure_filename
import tempfile
import logging

# Importações para extração de documentos
import docx
import pymupdf
import openpyxl
import pandas as pd
from PIL import Image
import zipfile
import json

# Configurar logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

extractor_bp = Blueprint('extractor', __name__)

# Configurações de upload expandidas
ALLOWED_EXTENSIONS = {
    'txt': 'text/plain',
    'pdf': 'application/pdf',
    'doc': 'application/msword',
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'xls': 'application/vnd.ms-excel',
    'csv': 'text/csv',
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'bmp': 'image/bmp',
    'tiff': 'image/tiff',
    'rtf': 'application/rtf'
}

# Tamanho máximo por tipo de arquivo (em bytes)
MAX_FILE_SIZES = {
    'pdf': 50 * 1024 * 1024,  # 50MB
    'docx': 20 * 1024 * 1024,  # 20MB
    'pptx': 25 * 1024 * 1024,  # 25MB
    'xlsx': 30 * 1024 * 1024,  # 30MB
    'default': 10 * 1024 * 1024  # 10MB
}

def allowed_file(filename):
    """Verifica se o arquivo é permitido"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def validate_url(url):
    """Valida se a URL é válida e acessível"""
    try:
        # Verificar se é uma URL válida
        if not validators.url(url):
            return False, "URL inválida"
        
        # Verificar se o protocolo é seguro
        parsed = urlparse(url)
        if parsed.scheme not in ['http', 'https']:
            return False, "Apenas URLs HTTP/HTTPS são permitidas"
        
        return True, "URL válida"
    except Exception as e:
        return False, f"Erro ao validar URL: {str(e)}"

def get_filename_from_url(url, content_disposition=None):
    """Extrai o nome do arquivo da URL ou do cabeçalho Content-Disposition"""
    try:
        # Tentar obter do Content-Disposition primeiro
        if content_disposition:
            import re
            cd_match = re.search(r'filename[*]?=([^;]+)', content_disposition)
            if cd_match:
                filename = cd_match.group(1).strip().strip('"\'')
                # Decodificar se necessário
                if 'filename*=' in content_disposition:
                    filename = unquote(filename.split("''")[-1])
                return secure_filename(filename)
        
        # Extrair da URL
        parsed = urlparse(url)
        filename = os.path.basename(parsed.path)
        
        if filename and '.' in filename:
            return secure_filename(filename)
        
        # Se não conseguir extrair, usar nome genérico baseado na extensão
        return "downloaded_document"
    except Exception:
        return "downloaded_document"

def download_file_from_url(url, max_size=50*1024*1024):
    """Baixa arquivo de uma URL com validações de segurança"""
    try:
        # Garantir que max_size seja um inteiro
        try:
            max_size = int(max_size)
        except (ValueError, TypeError):
            max_size = 50 * 1024 * 1024  # Valor padrão
        
        # Validar URL
        is_valid, message = validate_url(url)
        if not is_valid:
            raise Exception(message)
        
        # Headers seguros para a requisição
        headers = {
            'User-Agent': 'Document-Extractor-API/1.0 (File Processing Bot)',
            'Accept': '*/*',
            'Accept-Encoding': 'gzip, deflate',
            'Connection': 'keep-alive'
        }
        
        # Fazer requisição HEAD primeiro para verificar tamanho
        try:
            head_response = requests.head(url, headers=headers, timeout=10, allow_redirects=True)
            
            # Verificar tamanho do arquivo
            content_length = head_response.headers.get('content-length')
            if content_length and int(content_length) > max_size:
                raise Exception(f"Arquivo muito grande. Máximo permitido: {max_size // (1024*1024)}MB")
                
        except requests.exceptions.RequestException:
            # Se HEAD falhar, continuar com GET mas com stream
            pass
        
        # Fazer download do arquivo
        response = requests.get(url, headers=headers, timeout=30, stream=True, allow_redirects=True)
        response.raise_for_status()
        
        # Verificar Content-Type se disponível
        content_type = response.headers.get('content-type', '').lower()
        
        # Obter nome do arquivo
        content_disposition = response.headers.get('content-disposition')
        filename = get_filename_from_url(url, content_disposition)
        
        # Baixar arquivo em chunks para controlar o tamanho
        file_content = io.BytesIO()
        downloaded_size = 0
        
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                downloaded_size += len(chunk)
                if downloaded_size > max_size:
                    raise Exception(f"Arquivo muito grande. Máximo permitido: {max_size // (1024*1024)}MB")
                file_content.write(chunk)
        
        file_content.seek(0)
        
        # Tentar detectar extensão pelo content-type se não tiver no filename
        if '.' not in filename:
            extension = None
            if 'pdf' in content_type:
                extension = 'pdf'
            elif 'msword' in content_type or 'officedocument.wordprocessing' in content_type:
                extension = 'docx' if 'openxmlformats' in content_type else 'doc'
            elif 'spreadsheet' in content_type or 'excel' in content_type:
                extension = 'xlsx' if 'openxmlformats' in content_type else 'xls'
            elif 'text/plain' in content_type:
                extension = 'txt'
            elif 'text/csv' in content_type:
                extension = 'csv'
            elif 'image' in content_type:
                if 'jpeg' in content_type or 'jpg' in content_type:
                    extension = 'jpg'
                elif 'png' in content_type:
                    extension = 'png'
            
            if extension:
                filename = f"{filename}.{extension}"
        
        # Para arquivos de texto sem extensão (como LICENSE), tentar detectar pelo content-type
        elif not any(filename.lower().endswith(ext) for ext in ALLOWED_EXTENSIONS.keys()):
            if 'text/plain' in content_type or 'text/' in content_type:
                # Se é texto mas não tem extensão suportada, forçar .txt
                filename = f"{filename}.txt"
        
        return {
            'content': file_content,
            'filename': filename,
            'size': downloaded_size,
            'content_type': content_type,
            'url': url
        }
        
    except requests.exceptions.Timeout:
        raise Exception("Timeout ao baixar arquivo. Tente novamente.")
    except requests.exceptions.ConnectionError:
        raise Exception("Erro de conexão ao baixar arquivo.")
    except requests.exceptions.HTTPError as e:
        raise Exception(f"Erro HTTP ao baixar arquivo: {e.response.status_code}")
    except Exception as e:
        raise Exception(f"Erro ao baixar arquivo: {str(e)}")

def validate_file_size(file_size, file_extension):
    """Valida o tamanho do arquivo"""
    max_size = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
    return file_size <= max_size

def extract_text_from_docx(file_path):
    """Extrai texto, imagens e metadados de documentos Word (.docx)"""
    try:
        doc = docx.Document(file_path)
        text_content = []
        images = []
        
        # Extrair texto dos parágrafos
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Incluir informações de formatação básica
                runs_info = []
                for run in paragraph.runs:
                    if run.text.strip():
                        run_data = {
                            'text': run.text,
                            'bold': run.bold,
                            'italic': run.italic,
                            'underline': run.underline
                        }
                        runs_info.append(run_data)
                
                text_content.append({
                    'type': 'paragraph',
                    'text': paragraph.text,
                    'runs': runs_info
                })
        
        # Extrair texto de tabelas
        tables_data = []
        for table_idx, table in enumerate(doc.tables):
            table_rows = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_cells.append(cell.text.strip())
                if row_cells:
                    table_rows.append(row_cells)
            
            if table_rows:
                tables_data.append({
                    'table_index': table_idx + 1,
                    'rows': table_rows,
                    'text': '\n'.join([' | '.join(row) for row in table_rows])
                })
        
        # Extrair imagens do documento
        try:
            # Acessar o arquivo ZIP do documento Word
            import zipfile
            from PIL import Image
            import io
            
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                # Listar arquivos de mídia
                media_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                
                for idx, media_file in enumerate(media_files):
                    try:
                        # Ler dados da imagem
                        img_data = docx_zip.read(media_file)
                        
                        # Verificar se é imagem válida
                        try:
                            img = Image.open(io.BytesIO(img_data))
                            width, height = img.size
                            img_format = img.format or 'UNKNOWN'
                            
                            # Converter para base64
                            img_base64 = base64.b64encode(img_data).decode('utf-8')
                            
                            # Determinar tipo MIME
                            mime_type = f"image/{img_format.lower()}"
                            if img_format.upper() == 'JPEG':
                                mime_type = "image/jpeg"
                            elif img_format.upper() == 'PNG':
                                mime_type = "image/png"
                            
                            images.append({
                                'index': idx + 1,
                                'filename': media_file.split('/')[-1],
                                'format': img_format.upper(),
                                'width': width,
                                'height': height,
                                'size_bytes': len(img_data),
                                'data': img_base64,
                                'mime_type': mime_type
                            })
                            
                            logger.debug(f"Imagem extraída: {media_file} - {width}x{height} - {len(img_data)} bytes")
                            
                        except Exception as img_process_error:
                            logger.warning(f"Arquivo de mídia não é imagem válida: {media_file} - {img_process_error}")
                            
                    except Exception as media_error:
                        logger.warning(f"Erro ao processar mídia {media_file}: {media_error}")
                        
        except Exception as image_error:
            logger.warning(f"Erro ao extrair imagens do Word: {image_error}")
        
        # Extrair propriedades do documento
        core_properties = doc.core_properties
        metadata = {
            'title': core_properties.title or '',
            'author': core_properties.author or '',
            'subject': core_properties.subject or '',
            'created': core_properties.created.isoformat() if core_properties.created else None,
            'modified': core_properties.modified.isoformat() if core_properties.modified else None
        }
        
        # Texto combinado
        combined_text = '\n'.join([item['text'] for item in text_content])
        for table in tables_data:
            combined_text += '\n\n' + table['text']
        
        # Adicionar informação sobre imagens no texto
        if images:
            combined_text += f"\n\n[INFO: Documento contém {len(images)} imagem(s) extraída(s)]"
        
        return {
            'text': combined_text,
            'paragraphs': text_content,
            'tables': tables_data,
            'images': images,
            'metadata': metadata,
            'stats': {
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'image_count': len(images),
                'character_count': len(combined_text),
                'word_count': len(combined_text.split())
            }
        }
    except Exception as e:
        logger.error(f"Erro ao extrair dados do Word: {str(e)}")
        raise Exception(f"Erro ao extrair dados do Word: {str(e)}")

def extract_text_from_pptx(file_path):
    """Extrai texto, imagens e metadados de apresentações PowerPoint (.pptx)"""
    try:
        from pptx import Presentation
        import zipfile
        from PIL import Image
        import io
        
        prs = Presentation(file_path)
        text_content = []
        images = []
        slides_data = []
        
        # Extrair texto e layout de cada slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            slide_shapes = []
            
            # Processar formas no slide
            for shape in slide.shapes:
                shape_info = {
                    'type': str(shape.shape_type),
                    'text': ''
                }
                
                # Extrair texto de formas com texto
                if hasattr(shape, "text") and shape.text.strip():
                    shape_info['text'] = shape.text.strip()
                    slide_text.append(shape.text.strip())
                
                # Verificar se é uma tabela
                if hasattr(shape, "table"):
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(' | '.join(row_text))
                    
                    if table_text:
                        table_content = '\n'.join(table_text)
                        shape_info['text'] = table_content
                        slide_text.append(table_content)
                
                if shape_info['text']:
                    slide_shapes.append(shape_info)
            
            # Dados do slide
            slide_data = {
                'slide_number': slide_idx + 1,
                'text': '\n'.join(slide_text),
                'shapes': slide_shapes,
                'shape_count': len(slide.shapes)
            }
            
            slides_data.append(slide_data)
            
            if slide_text:
                text_content.append({
                    'slide': slide_idx + 1,
                    'text': '\n'.join(slide_text)
                })
        
        # Extrair imagens do arquivo PPTX
        try:
            with zipfile.ZipFile(file_path, 'r') as pptx_zip:
                # Listar arquivos de mídia
                media_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/media/')]
                
                for idx, media_file in enumerate(media_files):
                    try:
                        # Ler dados da imagem
                        img_data = pptx_zip.read(media_file)
                        
                        # Verificar se é imagem válida
                        try:
                            img = Image.open(io.BytesIO(img_data))
                            width, height = img.size
                            img_format = img.format or 'UNKNOWN'
                            
                            # Converter para base64
                            img_base64 = base64.b64encode(img_data).decode('utf-8')
                            
                            # Determinar tipo MIME
                            mime_type = f"image/{img_format.lower()}"
                            if img_format.upper() == 'JPEG':
                                mime_type = "image/jpeg"
                            elif img_format.upper() == 'PNG':
                                mime_type = "image/png"
                            
                            images.append({
                                'index': idx + 1,
                                'filename': media_file.split('/')[-1],
                                'format': img_format.upper(),
                                'width': width,
                                'height': height,
                                'size_bytes': len(img_data),
                                'data': img_base64,
                                'mime_type': mime_type
                            })
                            
                            logger.debug(f"Imagem extraída do PPTX: {media_file} - {width}x{height} - {len(img_data)} bytes")
                            
                        except Exception as img_process_error:
                            logger.warning(f"Arquivo de mídia não é imagem válida: {media_file} - {img_process_error}")
                            
                    except Exception as media_error:
                        logger.warning(f"Erro ao processar mídia {media_file}: {media_error}")
                        
        except Exception as image_error:
            logger.warning(f"Erro ao extrair imagens do PowerPoint: {image_error}")
        
        # Metadados da apresentação
        metadata = {
            'slide_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
        }
        
        # Tentar extrair propriedades do core
        try:
            core_props = prs.core_properties
            metadata.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': core_props.created.isoformat() if core_props.created else None,
                'modified': core_props.modified.isoformat() if core_props.modified else None
            })
        except Exception as props_error:
            logger.warning(f"Erro ao extrair propriedades: {props_error}")
        
        # Texto combinado
        combined_text = ''
        for slide_data in slides_data:
            if slide_data['text']:
                combined_text += f"--- Slide {slide_data['slide_number']} ---\n{slide_data['text']}\n\n"
        
        # Adicionar informação sobre imagens
        if images:
            combined_text += f"[INFO: Apresentação contém {len(images)} imagem(s) extraída(s)]"
        
        return {
            'text': combined_text.strip(),
            'slides': slides_data,
            'images': images,
            'metadata': metadata,
            'stats': {
                'slide_count': len(prs.slides),
                'image_count': len(images),
                'character_count': len(combined_text),
                'word_count': len(combined_text.split()),
                'total_shapes': sum(slide['shape_count'] for slide in slides_data)
            }
        }
        
    except Exception as e:
        logger.error(f"Erro ao extrair dados do PowerPoint: {str(e)}")
        raise Exception(f"Erro ao extrair dados do PowerPoint: {str(e)}")

def is_scanned_pdf(doc):
    """
    Detecta se um PDF é digitalizado/escaneado usando heurísticas avançadas
    
    Baseado nas melhores práticas de detecção:
    1. Verifica se páginas são dominadas por imagens
    2. Analisa a quantidade de texto extraível
    3. Detecta fontes específicas de OCR
    4. Calcula ratio de cobertura de imagens
    """
    try:
        total_pages = len(doc)
        scanned_indicators = 0
        total_text_chars = 0
        pages_with_large_images = 0
        ocr_fonts_detected = 0
        
        # Fontes comuns de OCR
        ocr_fonts = {
            'GlyphlessFont',  # Tesseract
            'Arial-BoldMT', 'ArialMT',  # Comum em OCR
            'TimesNewRomanPSMT',
            'CourierNewPSMT'
        }
        
        for page_num in range(total_pages):
            page = doc[page_num]
            page_rect = page.rect
            page_area = abs(page_rect)
            
            # 1. Verificar cobertura de imagens
            try:
                image_coverage = 0
                image_list = page.get_images()
                
                for img in image_list:
                    try:
                        # Obter bbox da imagem
                        img_dict = page.get_image_bbox(img[0])
                        if img_dict:
                            img_area = abs(img_dict)
                            coverage = img_area / page_area if page_area > 0 else 0
                            image_coverage += coverage
                    except:
                        continue
                
                # Se >90% da página é coberta por imagens, é provável que seja escaneada
                if image_coverage >= 0.9:
                    pages_with_large_images += 1
                    scanned_indicators += 2
                    
            except Exception as e:
                logger.debug(f"Erro ao analisar imagens na página {page_num}: {e}")
            
            # 2. Analisar texto e fontes
            try:
                text = page.get_text().strip()
                total_text_chars += len(text)
                
                # Verificar fontes OCR
                blocks = page.get_text("dict")
                page_fonts = set()
                
                for block in blocks.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line.get("spans", []):
                                font_name = span.get('font', '')
                                page_fonts.add(font_name)
                                
                                # Detectar fontes de OCR
                                if any(ocr_font in font_name for ocr_font in ocr_fonts):
                                    ocr_fonts_detected += 1
                                    scanned_indicators += 1
                
                # Pouco texto extraível indica PDF escaneado
                if len(text) < 50 and len(image_list) > 0:
                    scanned_indicators += 2
                    
            except Exception as e:
                logger.debug(f"Erro ao analisar texto na página {page_num}: {e}")
        
        # 3. Calcular score de confiança
        avg_text_per_page = total_text_chars / total_pages if total_pages > 0 else 0
        image_ratio = pages_with_large_images / total_pages if total_pages > 0 else 0
        
        # Critérios de decisão
        is_scanned = False
        confidence = 0
        
        if image_ratio >= 0.7:  # 70% das páginas dominadas por imagens
            is_scanned = True
            confidence = 0.8 + (image_ratio * 0.2)
        elif avg_text_per_page < 100 and pages_with_large_images > 0:  # Pouco texto + imagens
            is_scanned = True
            confidence = 0.7
        elif ocr_fonts_detected >= total_pages:  # Fontes OCR detectadas
            is_scanned = True
            confidence = 0.6
        elif scanned_indicators >= total_pages * 2:  # Score alto de indicadores
            is_scanned = True
            confidence = 0.5
        
        logger.info(f"Detecção PDF escaneado - Páginas: {total_pages}, Imagens dominantes: {pages_with_large_images}, "
                   f"Texto médio/página: {avg_text_per_page:.1f}, Fontes OCR: {ocr_fonts_detected}, "
                   f"Score: {scanned_indicators}, É escaneado: {is_scanned}, Confiança: {confidence:.2f}")
        
        return is_scanned, confidence
        
    except Exception as e:
        logger.error(f"Erro na detecção de PDF escaneado: {e}")
        return False, 0.0

def pdf_pages_to_images(doc, dpi=200, image_format='PNG'):
    """
    Converte páginas do PDF em imagens base64 de alta qualidade
    
    Args:
        doc: Documento PyMuPDF
        dpi: Resolução das imagens (default: 200 DPI para boa qualidade)
        image_format: Formato de saída ('PNG' ou 'JPEG')
    
    Returns:
        Lista de dicionários com informações das imagens
    """
    try:
        images = []
        
        # Configurar matriz de zoom baseada no DPI
        zoom = dpi / 72.0  # 72 DPI é o padrão
        mat = pymupdf.Matrix(zoom, zoom)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            try:
                # Renderizar página como imagem
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # Converter para bytes
                if image_format.upper() == 'JPEG':
                    img_data = pix.tobytes("jpeg", jpg_quality=95)
                    mime_type = "image/jpeg"
                else:
                    img_data = pix.tobytes("png")
                    mime_type = "image/png"
                
                # Converter para base64
                img_base64 = base64.b64encode(img_data).decode('utf-8')
                
                # Informações da imagem
                image_info = {
                    'page': page_num + 1,
                    'width': pix.width,
                    'height': pix.height,
                    'dpi': dpi,
                    'format': image_format.upper(),
                    'size_bytes': len(img_data),
                    'data': img_base64,
                    'mime_type': mime_type
                }
                
                images.append(image_info)
                pix = None  # Liberar memória
                
                logger.debug(f"Página {page_num + 1} convertida: {pix.width}x{pix.height} - {len(img_data)} bytes")
                
            except Exception as page_error:
                logger.error(f"Erro ao converter página {page_num + 1}: {page_error}")
                # Adicionar placeholder em caso de erro
                images.append({
                    'page': page_num + 1,
                    'error': str(page_error),
                    'width': 0,
                    'height': 0,
                    'data': None
                })
        
        return images
        
    except Exception as e:
        logger.error(f"Erro ao converter PDF para imagens: {e}")
        raise Exception(f"Erro ao converter PDF para imagens: {e}")

def extract_text_from_pdf(file_path):
    """Extrai texto, imagens e metadados de documentos PDF com fallback para PDFs escaneados"""
    try:
        text_content = []
        images = []
        fonts_info = []
        metadata = {}
        total_pages = 0
        is_scanned = False
        scanned_confidence = 0.0
        fallback_images = []
        
        # Usar context manager para garantir fechamento correto
        with pymupdf.open(file_path) as doc:
            # Extrair metadados
            metadata = {
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', ''),
                'creator': doc.metadata.get('creator', ''),
                'producer': doc.metadata.get('producer', ''),
                'creation_date': doc.metadata.get('creationDate', ''),
                'modification_date': doc.metadata.get('modDate', ''),
                'pages': len(doc)
            }
            
            total_pages = len(doc)
            
            # DETECÇÃO DE PDF ESCANEADO
            is_scanned, scanned_confidence = is_scanned_pdf(doc)
            
            # Processamento normal de extração
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extrair texto
                text = page.get_text()
                if text.strip():
                    text_content.append({
                        'page': page_num + 1,
                        'text': text.strip()
                    })
                
                # Extrair informações de fontes
                try:
                    blocks = page.get_text("dict")
                    page_fonts = set()
                    for block in blocks.get("blocks", []):
                        if "lines" in block:
                            for line in block["lines"]:
                                for span in line.get("spans", []):
                                    font_info = f"{span.get('font', 'Unknown')} - {span.get('size', 'Unknown')}pt"
                                    page_fonts.add(font_info)
                    
                    if page_fonts:
                        fonts_info.append({
                            'page': page_num + 1,
                            'fonts': list(page_fonts)
                        })
                except Exception as font_error:
                    logger.warning(f"Erro ao extrair fontes da página {page_num}: {font_error}")
                
                # Extrair imagens embutidas com melhor tratamento
                try:
                    image_list = page.get_images()
                    logger.debug(f"Página {page_num + 1}: encontradas {len(image_list)} imagens embutidas")
                    
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            
                            if base_image:
                                img_data = base_image["image"]
                                img_ext = base_image["ext"]
                                img_base64 = base64.b64encode(img_data).decode()
                                
                                # Obter dimensões usando Pixmap como fallback
                                try:
                                    pix = pymupdf.Pixmap(doc, xref)
                                    width, height = pix.width, pix.height
                                    pix = None
                                except:
                                    width, height = 0, 0
                                
                                # Determinar tipo MIME
                                mime_type = f"image/{img_ext}"
                                if img_ext.lower() in ['jpg', 'jpeg']:
                                    mime_type = "image/jpeg"
                                elif img_ext.lower() == 'png':
                                    mime_type = "image/png"
                                
                                images.append({
                                    'page': page_num + 1,
                                    'index': img_index + 1,
                                    'format': img_ext.upper(),
                                    'data': img_base64,
                                    'width': width,
                                    'height': height,
                                    'size_bytes': len(img_data),
                                    'mime_type': mime_type,
                                    'xref': xref
                                })
                                
                                logger.debug(f"Imagem extraída: página {page_num + 1}, índice {img_index + 1}, {width}x{height}, {len(img_data)} bytes")
                            
                        except Exception as img_error:
                            logger.warning(f"Erro ao extrair imagem {img_index} da página {page_num + 1}: {img_error}")
                            # Tentar método alternativo com Pixmap
                            try:
                                xref = img[0]
                                pix = pymupdf.Pixmap(doc, xref)
                                
                                if pix.n - pix.alpha < 4:  # GRAY ou RGB
                                    img_data = pix.tobytes("png")
                                    img_base64 = base64.b64encode(img_data).decode()
                                    
                                    images.append({
                                        'page': page_num + 1,
                                        'index': img_index + 1,
                                        'format': 'PNG',
                                        'data': img_base64,
                                        'width': pix.width,
                                        'height': pix.height,
                                        'size_bytes': len(img_data),
                                        'mime_type': 'image/png',
                                        'extracted_method': 'pixmap_fallback'
                                    })
                                    
                                    logger.debug(f"Imagem extraída (fallback): página {page_num + 1}, {pix.width}x{pix.height}")
                                
                                pix = None
                            except Exception as fallback_error:
                                logger.warning(f"Falha também no método alternativo para imagem {img_index}: {fallback_error}")
                                
                except Exception as page_error:
                    logger.warning(f"Erro ao processar imagens da página {page_num + 1}: {page_error}")
            
            # EXTRAÇÃO DE IMAGENS DE PÁGINAS (sempre ativar para visualização)
            page_images = []
            try:
                # Configurar qualidade baseada no número de imagens embutidas
                if len(images) == 0 or is_scanned:
                    # Se não há imagens embutidas ou é escaneado, converter páginas
                    if scanned_confidence >= 0.8:
                        dpi = 300  # Alta qualidade para PDFs claramente escaneados
                        image_format = 'PNG'  # PNG para preservar qualidade
                    elif scanned_confidence >= 0.5:
                        dpi = 250  # Qualidade média-alta
                        image_format = 'PNG'
                    else:
                        dpi = 150  # Qualidade padrão para visualização
                        image_format = 'JPEG'
                    
                    logger.info(f"Extraindo imagens de páginas com DPI {dpi} formato {image_format}")
                    page_images = pdf_pages_to_images(doc, dpi=dpi, image_format=image_format)
                    
                    if is_scanned and scanned_confidence >= 0.5:
                        logger.info(f"PDF detectado como escaneado (confiança: {scanned_confidence:.2f}). Páginas convertidas para preservar conteúdo.")
                        
                        # Criar texto indicativo do fallback para PDFs escaneados
                        if len(text_content) == 0:
                            fallback_text = f"[PDF DIGITALIZADO DETECTADO - {total_pages} páginas convertidas em imagens]\n\n"
                            fallback_text += f"Este documento foi identificado como digitalizado/escaneado.\n"
                            fallback_text += f"Confiança na detecção: {scanned_confidence:.1%}\n"
                            fallback_text += f"As páginas foram convertidas em imagens de {dpi} DPI para preservar o conteúdo.\n\n"
                            
                            for i, img in enumerate(page_images):
                                if 'error' not in img:
                                    fallback_text += f"--- Página {img['page']} (Imagem {img['width']}x{img['height']}) ---\n"
                                    fallback_text += f"[Imagem disponível em formato {img['format']} - {img['size_bytes']} bytes]\n\n"
                            
                            combined_text = fallback_text
                        else:
                            # Se há algum texto, adicionar aviso sobre fallback
                            combined_text = '\n\n'.join([f"--- Página {item['page']} ---\n{item['text']}" for item in text_content])
                            combined_text += f"\n\n[AVISO: PDF detectado como possivelmente escaneado (confiança: {scanned_confidence:.1%})."
                            combined_text += f" Imagens de páginas foram geradas para garantir preservação do conteúdo.]"
                    else:
                        # PDF normal com imagens de páginas para visualização
                        combined_text = '\n\n'.join([f"--- Página {item['page']} ---\n{item['text']}" for item in text_content])
                        if page_images:
                            combined_text += f"\n\n[INFO: {len(page_images)} páginas convertidas em imagens para visualização.]"
                else:
                    # PDF normal com imagens embutidas - só texto
                    combined_text = '\n\n'.join([f"--- Página {item['page']} ---\n{item['text']}" for item in text_content])
                    
            except Exception as page_img_error:
                logger.error(f"Erro ao extrair imagens de páginas: {page_img_error}")
                # Continuar com extração normal em caso de erro
                page_images = []
                combined_text = '\n\n'.join([f"--- Página {item['page']} ---\n{item['text']}" for item in text_content])
        
        # Preparar resultado final
        result = {
            'text': combined_text,
            'pages_content': text_content,
            'images': images,
            'fonts': fonts_info,
            'metadata': metadata,
            'stats': {
                'page_count': total_pages,
                'image_count': len(images),
                'character_count': len(combined_text),
                'word_count': len(combined_text.split()),
                'is_scanned': is_scanned,
                'scanned_confidence': scanned_confidence
            }
        }
        
        # Adicionar imagens de páginas se disponíveis
        if page_images:
            result['page_images'] = page_images
            result['stats']['page_images_count'] = len(page_images)
            result['stats']['total_page_images_size'] = sum(img.get('size_bytes', 0) for img in page_images if isinstance(img, dict))
            
            # Calcular tamanho total das imagens de páginas em MB
            total_size_mb = result['stats']['total_page_images_size'] / (1024 * 1024)
            result['stats']['total_page_images_size_mb'] = round(total_size_mb, 2)
            
            # Indicar se são imagens de fallback (para PDFs escaneados)
            if is_scanned and scanned_confidence >= 0.5:
                result['stats']['is_fallback'] = True
                result['fallback_images'] = page_images  # Manter compatibilidade
            else:
                result['stats']['is_fallback'] = False
        
        return result
    except Exception as e:
        logger.error(f"Erro ao extrair dados do PDF: {str(e)}")
        raise Exception(f"Erro ao extrair dados do PDF: {str(e)}")

def extract_data_from_excel(file_path):
    """Extrai dados de planilhas Excel com análise avançada"""
    try:
        # Usar context manager para garantir fechamento correto
        with pd.ExcelFile(file_path) as excel_file:
            sheets_data = {}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Análise básica dos dados
                sheet_analysis = {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist(),
                    'data_types': {str(k): str(v) for k, v in df.dtypes.to_dict().items()},
                    'null_counts': df.isnull().sum().to_dict()
                }
                
                # Estatísticas para colunas numéricas
                numeric_columns = df.select_dtypes(include=['number']).columns
                if len(numeric_columns) > 0:
                    sheet_analysis['numeric_stats'] = df[numeric_columns].describe().to_dict()
                
                # Converter DataFrame para texto estruturado
                text_content = f"=== Planilha: {sheet_name} ===\n"
                text_content += f"Linhas: {len(df)}, Colunas: {len(df.columns)}\n"
                text_content += f"Colunas: {', '.join(df.columns.tolist())}\n\n"
                text_content += df.to_string(index=False, max_rows=1000)  # Limitar a 1000 linhas
                
                sheets_data[sheet_name] = {
                    'text': text_content,
                    'analysis': sheet_analysis,
                    'sample_data': df.head(10).to_dict('records') if len(df) > 0 else []
                }
            
            # Texto combinado de todas as planilhas
            combined_text = '\n\n'.join([data['text'] for data in sheets_data.values()])
            total_sheets = len(excel_file.sheet_names)
            
            return {
                'text': combined_text,
                'sheets': sheets_data,
                'stats': {
                    'total_sheets': total_sheets,
                    'total_rows': sum([data['analysis']['rows'] for data in sheets_data.values()]),
                    'total_columns': sum([data['analysis']['columns'] for data in sheets_data.values()]),
                    'character_count': len(combined_text),
                    'word_count': len(combined_text.split())
                }
            }
    except Exception as e:
        logger.error(f"Erro ao extrair dados do Excel: {str(e)}")
        raise Exception(f"Erro ao extrair dados do Excel: {str(e)}")

def extract_text_from_csv(file_path):
    """Extrai dados de arquivos CSV"""
    try:
        # Tentar diferentes encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        df = None
        used_encoding = None
        
        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue
        
        if df is None:
            raise Exception("Não foi possível decodificar o arquivo CSV")
        
        # Análise dos dados
        analysis = {
            'rows': len(df),
            'columns': len(df.columns),
            'column_names': df.columns.tolist(),
            'data_types': {str(k): str(v) for k, v in df.dtypes.to_dict().items()},
            'null_counts': df.isnull().sum().to_dict(),
            'encoding_used': used_encoding
        }
        
        # Estatísticas para colunas numéricas
        numeric_columns = df.select_dtypes(include=['number']).columns
        if len(numeric_columns) > 0:
            analysis['numeric_stats'] = df[numeric_columns].describe().to_dict()
        
        # Converter para texto
        text_content = f"=== Arquivo CSV ===\n"
        text_content += f"Linhas: {len(df)}, Colunas: {len(df.columns)}\n"
        text_content += f"Encoding: {used_encoding}\n"
        text_content += f"Colunas: {', '.join(df.columns.tolist())}\n\n"
        text_content += df.to_string(index=False, max_rows=1000)
        
        return {
            'text': text_content,
            'analysis': analysis,
            'sample_data': df.head(10).to_dict('records') if len(df) > 0 else [],
            'stats': {
                'rows': len(df),
                'columns': len(df.columns),
                'character_count': len(text_content),
                'word_count': len(text_content.split()),
                'encoding_used': used_encoding
            }
        }
    except Exception as e:
        logger.error(f"Erro ao extrair dados do CSV: {str(e)}")
        raise Exception(f"Erro ao extrair dados do CSV: {str(e)}")

def extract_text_from_txt(file_path):
    """Extrai texto de arquivos de texto simples com detecção de encoding"""
    try:
        # Tentar diferentes encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        content = None
        used_encoding = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                    used_encoding = encoding
                    break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            raise Exception("Não foi possível decodificar o arquivo de texto")
        
        lines = content.split('\n')
        
        # Análise do conteúdo
        analysis = {
            'lines': len(lines),
            'characters': len(content),
            'words': len(content.split()),
            'paragraphs': len([line for line in lines if line.strip()]),
            'encoding_used': used_encoding,
            'empty_lines': len([line for line in lines if not line.strip()])
        }
        
        return {
            'text': content,
            'analysis': analysis,
            'stats': {
                'lines': len(lines),
                'characters': len(content),
                'words': len(content.split()),
                'encoding_used': used_encoding
            }
        }
        
    except Exception as e:
        logger.error(f"Erro ao extrair texto: {str(e)}")
        raise Exception(f"Erro ao extrair texto: {str(e)}")

def extract_text_from_image(file_path):
    """Extrai metadados e dados de imagens em formato padronizado"""
    try:
        with Image.open(file_path) as img:
            # Metadados básicos
            metadata = {
                'format': img.format,
                'mode': img.mode,
                'size': img.size,
                'width': img.width,
                'height': img.height,
                'has_transparency': img.mode in ('RGBA', 'LA') or 'transparency' in img.info
            }
            
            # Ler arquivo original para base64
            with open(file_path, 'rb') as img_file:
                img_data = img_file.read()
                img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            # Determinar tipo MIME
            mime_type = f"image/{img.format.lower()}" if img.format else "image/unknown"
            if img.format and img.format.upper() == 'JPEG':
                mime_type = "image/jpeg"
            elif img.format and img.format.upper() == 'PNG':
                mime_type = "image/png"
            
            # Preparar dados da imagem no formato padrão
            image_data = {
                'index': 1,
                'filename': file_path.split('/')[-1] if '/' in file_path else file_path.split('\\')[-1],
                'format': img.format or 'UNKNOWN',
                'width': img.width,
                'height': img.height,
                'mode': img.mode,
                'size_bytes': len(img_data),
                'data': img_base64,
                'mime_type': mime_type
            }
            
            # Texto descritivo
            description = f"Imagem detectada: {img.format} {img.width}x{img.height} pixels, Modo: {img.mode}"
            ocr_note = "[OCR não configurado - Para extrair texto de imagens, configure o Tesseract OCR]"
            combined_text = f"{description}\n\n{ocr_note}\n\n[INFO: Imagem incluída em base64 para visualização]"
            
            return {
                'text': combined_text,
                'images': [image_data],
                'metadata': metadata,
                'stats': {
                    'format': img.format,
                    'width': img.width,
                    'height': img.height,
                    'mode': img.mode,
                    'image_count': 1,
                    'size_pixels': f"{img.width}x{img.height}",
                    'character_count': len(combined_text),
                    'word_count': len(combined_text.split())
                }
            }
    except Exception as e:
        logger.error(f"Erro ao processar imagem: {str(e)}")
        raise Exception(f"Erro ao processar imagem: {str(e)}")

@extractor_bp.route('/extract', methods=['POST'])
def extract_document():
    """Endpoint principal para extração de documentos com validações aprimoradas"""
    try:
        # Verificar se um arquivo foi enviado
        if 'file' not in request.files:
            return jsonify({
                'success': False,
                'error': 'Nenhum arquivo foi enviado',
                'error_code': 'NO_FILE'
            }), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({
                'success': False,
                'error': 'Nenhum arquivo selecionado',
                'error_code': 'EMPTY_FILENAME'
            }), 400
        
        # Validar tipo de arquivo
        if not allowed_file(file.filename):
            return jsonify({
                'success': False,
                'error': 'Tipo de arquivo não suportado',
                'error_code': 'UNSUPPORTED_TYPE',
                'supported_types': list(ALLOWED_EXTENSIONS.keys())
            }), 400
        
        filename = secure_filename(file.filename)
        file_extension = filename.rsplit('.', 1)[1].lower()
        
        # Validar tamanho do arquivo
        file.seek(0, 2)  # Mover para o final
        file_size = file.tell()
        file.seek(0)  # Voltar ao início
        
        if not validate_file_size(file_size, file_extension):
            max_size = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
            return jsonify({
                'success': False,
                'error': f'Arquivo muito grande. Tamanho máximo: {max_size // (1024*1024)}MB',
                'error_code': 'FILE_TOO_LARGE',
                'max_size_mb': max_size // (1024*1024)
            }), 413
        
        # Salvar arquivo temporariamente
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            file.save(temp_file.name)
            temp_file_path = temp_file.name
        
        try:
            # Processar arquivo baseado na extensão
            if file_extension == 'docx':
                result = extract_text_from_docx(temp_file_path)
            elif file_extension == 'pdf':
                result = extract_text_from_pdf(temp_file_path)
            elif file_extension in ['xlsx', 'xls']:
                result = extract_data_from_excel(temp_file_path)
            elif file_extension == 'csv':
                result = extract_text_from_csv(temp_file_path)
            elif file_extension == 'txt':
                result = extract_text_from_txt(temp_file_path)
            elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                result = extract_text_from_image(temp_file_path)
            else:
                return jsonify({
                    'success': False,
                    'error': f'Processamento para {file_extension} não implementado',
                    'error_code': 'PROCESSING_NOT_IMPLEMENTED'
                }), 400
            
            # Adicionar informações do arquivo
            result['file_info'] = {
                'filename': filename,
                'original_filename': file.filename,
                'type': file_extension,
                'mime_type': ALLOWED_EXTENSIONS.get(file_extension, 'unknown'),
                'size_bytes': file_size,
                'size_mb': round(file_size / (1024*1024), 2)
            }
            
            logger.info(f"Arquivo processado com sucesso: {filename} ({file_size} bytes)")
            
            return jsonify({
                'success': True,
                'data': result,
                'message': 'Documento processado com sucesso'
            })
        
        finally:
            # Limpar arquivo temporário
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
    
    except Exception as e:
        logger.error(f"Erro no processamento: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'error_code': 'PROCESSING_ERROR'
        }), 500

@extractor_bp.route('/extract/bulk', methods=['POST'])
def extract_multiple_documents():
    """Endpoint para processamento em lote de múltiplos documentos"""
    try:
        files = request.files.getlist('files')
        
        if not files or len(files) == 0:
            return jsonify({
                'success': False,
                'error': 'Nenhum arquivo foi enviado',
                'error_code': 'NO_FILES'
            }), 400
        
        if len(files) > 10:  # Limitar a 10 arquivos por vez
            return jsonify({
                'success': False,
                'error': 'Máximo de 10 arquivos por vez',
                'error_code': 'TOO_MANY_FILES'
            }), 400
        
        results = []
        errors = []
        total_processing_time = 0
        
        for i, file in enumerate(files):
            start_time = time.time()
            try:
                if file.filename == '':
                    errors.append({
                        'index': i,
                        'error': 'Nome do arquivo vazio',
                        'error_code': 'EMPTY_FILENAME'
                    })
                    continue
                
                # Verificar se o arquivo é permitido
                if not allowed_file(file.filename):
                    errors.append({
                        'index': i,
                        'filename': file.filename,
                        'error': 'Tipo de arquivo não suportado',
                        'error_code': 'UNSUPPORTED_TYPE'
                    })
                    continue
                
                filename = secure_filename(file.filename)
                file_extension = filename.rsplit('.', 1)[1].lower()
                
                # Validar tamanho do arquivo
                file.seek(0, 2)  # Mover para o final
                file_size = file.tell()
                file.seek(0)  # Voltar ao início
                
                if not validate_file_size(file_size, file_extension):
                    max_size = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
                    errors.append({
                        'index': i,
                        'filename': file.filename,
                        'error': f'Arquivo muito grande. Máximo: {max_size // (1024*1024)}MB',
                        'error_code': 'FILE_TOO_LARGE'
                    })
                    continue
                
                # Salvar arquivo temporariamente
                with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
                    file.save(temp_file.name)
                    temp_file_path = temp_file.name
                
                try:
                    # Processar arquivo baseado na extensão
                    if file_extension == 'docx':
                        result = extract_text_from_docx(temp_file_path)
                    elif file_extension == 'pdf':
                        result = extract_text_from_pdf(temp_file_path)
                    elif file_extension in ['xlsx', 'xls']:
                        result = extract_data_from_excel(temp_file_path)
                    elif file_extension == 'csv':
                        result = extract_text_from_csv(temp_file_path)
                    elif file_extension == 'txt':
                        result = extract_text_from_txt(temp_file_path)
                    elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                        result = extract_text_from_image(temp_file_path)
                    else:
                        errors.append({
                            'index': i,
                            'filename': file.filename,
                            'error': f'Processamento para {file_extension} não implementado',
                            'error_code': 'PROCESSING_NOT_IMPLEMENTED'
                        })
                        continue
                    
                    # Calcular tempo de processamento
                    processing_time = time.time() - start_time
                    total_processing_time += processing_time
                    
                    # Adicionar informações do arquivo
                    result['file_info'] = {
                        'filename': filename,
                        'original_filename': file.filename,
                        'type': file_extension,
                        'mime_type': ALLOWED_EXTENSIONS.get(file_extension, 'unknown'),
                        'size_bytes': file_size,
                        'size_mb': round(file_size / (1024*1024), 2),
                        'processing_time': round(processing_time, 2)
                    }
                    
                    results.append({
                        'index': i,
                        'filename': file.filename,
                        'success': True,
                        'data': result
                    })
                    
                    logger.info(f"Arquivo processado em lote: {filename} ({file_size} bytes) em {processing_time:.2f}s")
                
                finally:
                    # Limpar arquivo temporário
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                
            except Exception as e:
                processing_time = time.time() - start_time
                total_processing_time += processing_time
                
                logger.error(f"Erro ao processar arquivo {file.filename}: {str(e)}")
                errors.append({
                    'index': i,
                    'filename': file.filename,
                    'error': str(e),
                    'error_code': 'PROCESSING_ERROR',
                    'processing_time': round(processing_time, 2)
                })
        
        # Estatísticas finais
        total_files = len(files)
        processed_count = len(results)
        failed_count = len(errors)
        success_rate = (processed_count / total_files * 100) if total_files > 0 else 0
        
        return jsonify({
            'success': True,
            'results': results,
            'errors': errors,
            'summary': {
                'total_files': total_files,
                'processed': processed_count,
                'failed': failed_count,
                'success_rate': round(success_rate, 1),
                'total_processing_time': round(total_processing_time, 2),
                'average_time_per_file': round(total_processing_time / total_files, 2) if total_files > 0 else 0
            }
        })
        
    except Exception as e:
        logger.error(f"Erro no processamento em lote: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'error_code': 'BULK_PROCESSING_ERROR'
        }), 500

@extractor_bp.route('/supported-types', methods=['GET'])
def get_supported_types():
    """Retorna os tipos de arquivo suportados com informações detalhadas"""
    return jsonify({
        'success': True,
        'supported_extensions': list(ALLOWED_EXTENSIONS.keys()),
        'mime_types': ALLOWED_EXTENSIONS,
        'max_file_sizes': {
            ext: f"{size // (1024*1024)}MB" 
            for ext, size in MAX_FILE_SIZES.items()
        },
        'features_by_type': {
            'pdf': ['texto', 'imagens', 'metadados', 'fontes'],
            'docx': ['texto', 'formatação', 'tabelas', 'metadados'],
            'xlsx': ['dados', 'múltiplas planilhas', 'estatísticas', 'tipos de dados'],
            'csv': ['dados tabulares', 'detecção de encoding', 'estatísticas'],
            'txt': ['texto simples', 'detecção de encoding', 'análise básica'],
            'imagens': ['metadados', 'EXIF', 'dimensões', 'formato']
        }
    })

@extractor_bp.route('/health', methods=['GET'])
def health_check():
    """Endpoint de verificação de saúde da API"""
    import datetime
    return jsonify({
        'status': 'healthy',
        'service': 'Document Extractor API',
        'version': '2.0.0',
        'features': [
            'Extração de texto de PDF',
            'Extração de texto de Word',
            'Extração de dados de Excel',
            'Extração de dados de CSV',
            'Processamento de imagens',
            'Análise de metadados',
            'Validação de arquivos'
        ],
        'timestamp': datetime.datetime.now().isoformat()
    })

@extractor_bp.route('/stats', methods=['GET'])
def get_extraction_stats():
    """Endpoint para estatísticas de uso da API"""
    # Este endpoint pode ser expandido para incluir estatísticas reais
    return jsonify({
        'success': True,
        'stats': {
            'supported_formats': len(ALLOWED_EXTENSIONS),
            'max_file_size_mb': max(MAX_FILE_SIZES.values()) // (1024*1024),
            'features': {
                'text_extraction': True,
                'image_extraction': True,
                'metadata_extraction': True,
                'bulk_processing': True,
                'format_validation': True,
                'size_validation': True,
                'url_extraction': True
            }
        }
    })

@extractor_bp.route('/extract/url', methods=['POST'])
def extract_document_from_url():
    """Extrai dados de documento a partir de uma URL"""
    try:
        # Verificar se JSON foi enviado
        if not request.is_json:
            return jsonify({
                'success': False,
                'error': 'Content-Type deve ser application/json',
                'error_code': 'INVALID_CONTENT_TYPE'
            }), 400
        
        data = request.get_json()
        
        # Verificar se URL foi fornecida
        if not data or 'url' not in data:
            return jsonify({
                'success': False,
                'error': 'URL é obrigatória',
                'error_code': 'MISSING_URL'
            }), 400
        
        url = data['url'].strip()
        if not url:
            return jsonify({
                'success': False,
                'error': 'URL não pode estar vazia',
                'error_code': 'EMPTY_URL'
            }), 400
        
        # Obter tamanho máximo opcional (converter para int se necessário)
        max_size_mb = data.get('max_size_mb', 50)
        try:
            max_size_mb = int(float(max_size_mb))  # Suporta tanto int quanto float/string
        except (ValueError, TypeError):
            max_size_mb = 50  # Valor padrão em caso de erro
        max_size = max_size_mb * 1024 * 1024  # Converter MB para bytes
        
        logger.info(f"Iniciando download de: {url}")
        
        # Baixar arquivo da URL
        try:
            downloaded_file = download_file_from_url(url, max_size)
        except Exception as e:
            logger.error(f"Erro ao baixar arquivo de {url}: {str(e)}")
            return jsonify({
                'success': False,
                'error': str(e),
                'error_code': 'DOWNLOAD_ERROR',
                'url': url
            }), 400
        
        # Verificar se o arquivo baixado é suportado
        filename = downloaded_file['filename']
        if not allowed_file(filename):
            return jsonify({
                'success': False,
                'error': f'Tipo de arquivo não suportado: {filename}',
                'error_code': 'UNSUPPORTED_TYPE',
                'url': url,
                'filename': filename,
                'supported_types': list(ALLOWED_EXTENSIONS.keys())
            }), 400
        
        # Obter extensão do arquivo
        file_extension = filename.rsplit('.', 1)[1].lower()
        
        # Validar tamanho do arquivo baixado
        file_size = downloaded_file['size']
        if not validate_file_size(file_size, file_extension):
            max_allowed = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
            return jsonify({
                'success': False,
                'error': f'Arquivo muito grande. Máximo permitido: {max_allowed // (1024*1024)}MB',
                'error_code': 'FILE_TOO_LARGE',
                'url': url,
                'filename': filename,
                'file_size_mb': round(file_size / (1024*1024), 2)
            }), 400
        
        # Salvar arquivo temporariamente para processamento
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            temp_file.write(downloaded_file['content'].getvalue())
            temp_file_path = temp_file.name
        
        start_time = time.time()
        
        try:
            # Processar arquivo baseado na extensão
            if file_extension == 'docx':
                result = extract_text_from_docx(temp_file_path)
            elif file_extension == 'pdf':
                result = extract_text_from_pdf(temp_file_path)
            elif file_extension in ['xlsx', 'xls']:
                result = extract_data_from_excel(temp_file_path)
            elif file_extension == 'csv':
                result = extract_text_from_csv(temp_file_path)
            elif file_extension == 'txt':
                result = extract_text_from_txt(temp_file_path)
            elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                result = extract_text_from_image(temp_file_path)
            else:
                return jsonify({
                    'success': False,
                    'error': f'Processamento para {file_extension} não implementado',
                    'error_code': 'PROCESSING_NOT_IMPLEMENTED',
                    'url': url,
                    'filename': filename
                }), 400
            
            processing_time = time.time() - start_time
            
            # Adicionar informações do arquivo e download
            result['file_info'] = {
                'filename': filename,
                'type': file_extension,
                'mime_type': ALLOWED_EXTENSIONS.get(file_extension, 'unknown'),
                'size_bytes': file_size,
                'size_mb': round(file_size / (1024*1024), 2),
                'processing_time': round(processing_time, 2)
            }
            
            result['download_info'] = {
                'source_url': url,
                'content_type': downloaded_file['content_type'],
                'download_successful': True,
                'filename_from_url': filename
            }
            
            logger.info(f"Arquivo baixado e processado com sucesso: {filename} de {url} ({file_size} bytes) em {processing_time:.2f}s")
            
            return jsonify({
                'success': True,
                'data': result
            })
            
        finally:
            # Limpar arquivo temporário
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
    
    except Exception as e:
        logger.error(f"Erro na extração via URL: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'error_code': 'EXTRACTION_ERROR'
        }), 500

@extractor_bp.route('/extract/data', methods=['POST'])
def extract_document_from_data():
    """Endpoint para extração de documentos a partir de dados (base64 ou bytes)"""
    try:
        # Verificar se dados foram enviados
        if not request.json:
            return jsonify({
                'success': False,
                'error': 'Dados JSON são obrigatórios',
                'error_code': 'NO_JSON_DATA'
            }), 400
        
        data = request.json
        
        # Verificar campos obrigatórios
        if 'file_data' not in data:
            return jsonify({
                'success': False,
                'error': 'Campo "file_data" é obrigatório',
                'error_code': 'MISSING_FILE_DATA'
            }), 400
        
        if 'filename' not in data:
            return jsonify({
                'success': False,
                'error': 'Campo "filename" é obrigatório',
                'error_code': 'MISSING_FILENAME'
            }), 400
        
        filename = secure_filename(data['filename'])
        file_data = data['file_data']
        encoding = data.get('encoding', 'base64')  # base64 ou binary
        
        # Validar tipo de arquivo
        if not allowed_file(filename):
            return jsonify({
                'success': False,
                'error': 'Tipo de arquivo não suportado',
                'error_code': 'UNSUPPORTED_TYPE',
                'supported_types': list(ALLOWED_EXTENSIONS.keys())
            }), 400
        
        file_extension = filename.rsplit('.', 1)[1].lower()
        
        # Processar dados baseado na codificação
        try:
            if encoding == 'base64':
                # Remover prefixos data URI se presente
                if ',' in file_data:
                    file_data = file_data.split(',')[1]
                
                # Decodificar base64
                file_bytes = base64.b64decode(file_data)
            elif encoding == 'binary':
                # Assumir que os dados são bytes diretos
                if isinstance(file_data, str):
                    file_bytes = file_data.encode('latin-1')
                else:
                    file_bytes = file_data
            else:
                return jsonify({
                    'success': False,
                    'error': 'Encoding deve ser "base64" ou "binary"',
                    'error_code': 'INVALID_ENCODING'
                }), 400
                
        except Exception as e:
            return jsonify({
                'success': False,
                'error': f'Erro ao decodificar dados: {str(e)}',
                'error_code': 'DECODE_ERROR'
            }), 400
        
        file_size = len(file_bytes)
        
        # Validar tamanho do arquivo
        if not validate_file_size(file_size, file_extension):
            max_size = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
            return jsonify({
                'success': False,
                'error': f'Arquivo muito grande. Tamanho máximo: {max_size // (1024*1024)}MB',
                'error_code': 'FILE_TOO_LARGE',
                'max_size_mb': max_size // (1024*1024)
            }), 413
        
        # Salvar arquivo temporariamente
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        try:
            # Processar arquivo baseado na extensão
            if file_extension == 'docx':
                result = extract_text_from_docx(temp_file_path)
            elif file_extension == 'pdf':
                result = extract_text_from_pdf(temp_file_path)
            elif file_extension in ['xlsx', 'xls']:
                result = extract_data_from_excel(temp_file_path)
            elif file_extension == 'csv':
                result = extract_text_from_csv(temp_file_path)
            elif file_extension == 'txt':
                result = extract_text_from_txt(temp_file_path)
            elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                result = extract_text_from_image(temp_file_path)
            else:
                return jsonify({
                    'success': False,
                    'error': f'Processamento para {file_extension} não implementado',
                    'error_code': 'PROCESSING_NOT_IMPLEMENTED'
                }), 400
            
            # Adicionar informações do arquivo
            result['file_info'] = {
                'filename': filename,
                'type': file_extension,
                'mime_type': ALLOWED_EXTENSIONS.get(file_extension, 'unknown'),
                'size_bytes': file_size,
                'size_mb': round(file_size / (1024*1024), 2),
                'encoding_used': encoding
            }
            
            logger.info(f"Arquivo processado com sucesso via dados: {filename} ({file_size} bytes)")
            
            return jsonify({
                'success': True,
                'data': result,
                'message': 'Documento processado com sucesso a partir de dados'
            })
        
        finally:
            # Limpar arquivo temporário
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
    
    except Exception as e:
        logger.error(f"Erro no processamento de dados: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'error_code': 'PROCESSING_ERROR'
        }), 500

@extractor_bp.route('/extract/data/bulk', methods=['POST'])
def extract_multiple_documents_from_data():
    """Endpoint para processamento em lote de múltiplos documentos a partir de dados"""
    try:
        if not request.json:
            return jsonify({
                'success': False,
                'error': 'Dados JSON são obrigatórios',
                'error_code': 'NO_JSON_DATA'
            }), 400
        
        data = request.json
        
        if 'documents' not in data or not isinstance(data['documents'], list):
            return jsonify({
                'success': False,
                'error': 'Campo "documents" é obrigatório e deve ser uma lista',
                'error_code': 'INVALID_DOCUMENTS'
            }), 400
        
        documents = data['documents']
        
        if len(documents) == 0:
            return jsonify({
                'success': False,
                'error': 'Nenhum documento foi enviado',
                'error_code': 'NO_DOCUMENTS'
            }), 400
        
        if len(documents) > 10:  # Limitar a 10 documentos por vez
            return jsonify({
                'success': False,
                'error': 'Máximo de 10 documentos por vez',
                'error_code': 'TOO_MANY_DOCUMENTS'
            }), 400
        
        results = []
        errors = []
        total_processing_time = 0
        
        for i, doc in enumerate(documents):
            start_time = time.time()
            try:
                # Verificar campos obrigatórios
                if 'file_data' not in doc or 'filename' not in doc:
                    errors.append({
                        'index': i,
                        'error': 'Campos "file_data" e "filename" são obrigatórios',
                        'error_code': 'MISSING_REQUIRED_FIELDS'
                    })
                    continue
                
                filename = secure_filename(doc['filename'])
                file_data = doc['file_data']
                encoding = doc.get('encoding', 'base64')
                
                # Validar tipo de arquivo
                if not allowed_file(filename):
                    errors.append({
                        'index': i,
                        'filename': filename,
                        'error': 'Tipo de arquivo não suportado',
                        'error_code': 'UNSUPPORTED_TYPE'
                    })
                    continue
                
                file_extension = filename.rsplit('.', 1)[1].lower()
                
                # Processar dados baseado na codificação
                try:
                    if encoding == 'base64':
                        # Remover prefixos data URI se presente
                        if ',' in file_data:
                            file_data = file_data.split(',')[1]
                        file_bytes = base64.b64decode(file_data)
                    elif encoding == 'binary':
                        if isinstance(file_data, str):
                            file_bytes = file_data.encode('latin-1')
                        else:
                            file_bytes = file_data
                    else:
                        errors.append({
                            'index': i,
                            'filename': filename,
                            'error': 'Encoding deve ser "base64" ou "binary"',
                            'error_code': 'INVALID_ENCODING'
                        })
                        continue
                        
                except Exception as e:
                    errors.append({
                        'index': i,
                        'filename': filename,
                        'error': f'Erro ao decodificar dados: {str(e)}',
                        'error_code': 'DECODE_ERROR'
                    })
                    continue
                
                file_size = len(file_bytes)
                
                # Validar tamanho do arquivo
                if not validate_file_size(file_size, file_extension):
                    max_size = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
                    errors.append({
                        'index': i,
                        'filename': filename,
                        'error': f'Arquivo muito grande. Máximo: {max_size // (1024*1024)}MB',
                        'error_code': 'FILE_TOO_LARGE'
                    })
                    continue
                
                # Salvar arquivo temporariamente
                with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
                    temp_file.write(file_bytes)
                    temp_file_path = temp_file.name
                
                try:
                    # Processar arquivo baseado na extensão
                    if file_extension == 'docx':
                        result = extract_text_from_docx(temp_file_path)
                    elif file_extension == 'pdf':
                        result = extract_text_from_pdf(temp_file_path)
                    elif file_extension in ['xlsx', 'xls']:
                        result = extract_data_from_excel(temp_file_path)
                    elif file_extension == 'csv':
                        result = extract_text_from_csv(temp_file_path)
                    elif file_extension == 'txt':
                        result = extract_text_from_txt(temp_file_path)
                    elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                        result = extract_text_from_image(temp_file_path)
                    else:
                        errors.append({
                            'index': i,
                            'filename': filename,
                            'error': f'Processamento para {file_extension} não implementado',
                            'error_code': 'PROCESSING_NOT_IMPLEMENTED'
                        })
                        continue
                    
                    # Calcular tempo de processamento
                    processing_time = time.time() - start_time
                    total_processing_time += processing_time
                    
                    # Adicionar informações do arquivo
                    result['file_info'] = {
                        'filename': filename,
                        'type': file_extension,
                        'mime_type': ALLOWED_EXTENSIONS.get(file_extension, 'unknown'),
                        'size_bytes': file_size,
                        'size_mb': round(file_size / (1024*1024), 2),
                        'encoding_used': encoding,
                        'processing_time_seconds': round(processing_time, 3)
                    }
                    
                    results.append({
                        'index': i,
                        'success': True,
                        'data': result
                    })
                    
                finally:
                    # Limpar arquivo temporário
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                        
            except Exception as e:
                errors.append({
                    'index': i,
                    'filename': doc.get('filename', 'unknown'),
                    'error': str(e),
                    'error_code': 'PROCESSING_ERROR'
                })
        
        # Estatísticas
        total_documents = len(documents)
        successful_documents = len(results)
        failed_documents = len(errors)
        
        logger.info(f"Processamento em lote via dados: {successful_documents}/{total_documents} sucessos")
        
        return jsonify({
            'success': True,
            'data': {
                'results': results,
                'errors': errors,
                'statistics': {
                    'total_documents': total_documents,
                    'successful': successful_documents,
                    'failed': failed_documents,
                    'success_rate': round((successful_documents / total_documents) * 100, 2) if total_documents > 0 else 0,
                    'total_processing_time_seconds': round(total_processing_time, 3)
                }
            },
            'message': f'Processamento concluído: {successful_documents}/{total_documents} documentos processados com sucesso'
        })
    
    except Exception as e:
        logger.error(f"Erro no processamento em lote de dados: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'error_code': 'BULK_PROCESSING_ERROR'
        }), 500

@extractor_bp.route('/extract/url/bulk', methods=['POST'])
def extract_multiple_documents_from_urls():
    """Extrai dados de múltiplos documentos a partir de URLs"""
    try:
        # Verificar se JSON foi enviado
        if not request.is_json:
            return jsonify({
                'success': False,
                'error': 'Content-Type deve ser application/json',
                'error_code': 'INVALID_CONTENT_TYPE'
            }), 400
        
        data = request.get_json()
        
        # Verificar se URLs foram fornecidas
        if not data or 'urls' not in data:
            return jsonify({
                'success': False,
                'error': 'Lista de URLs é obrigatória',
                'error_code': 'MISSING_URLS'
            }), 400
        
        urls = data['urls']
        if not isinstance(urls, list) or not urls:
            return jsonify({
                'success': False,
                'error': 'URLs deve ser uma lista não vazia',
                'error_code': 'INVALID_URLS_FORMAT'
            }), 400
        
        # Limitar número de URLs
        max_urls = 5  # Limite para evitar sobrecarga
        if len(urls) > max_urls:
            return jsonify({
                'success': False,
                'error': f'Máximo de {max_urls} URLs permitidas por requisição',
                'error_code': 'TOO_MANY_URLS'
            }), 400
        
        # Obter configurações opcionais (converter para int se necessário)
        max_size_mb = data.get('max_size_mb', 50)
        try:
            max_size_mb = int(float(max_size_mb))  # Suporta tanto int quanto float/string
        except (ValueError, TypeError):
            max_size_mb = 50  # Valor padrão em caso de erro
        max_size = max_size_mb * 1024 * 1024
        
        results = []
        errors = []
        total_processing_time = 0
        
        for i, url in enumerate(urls):
            start_time = time.time()
            try:
                if not url or not isinstance(url, str):
                    errors.append({
                        'index': i,
                        'url': url,
                        'error': 'URL inválida ou vazia',
                        'error_code': 'INVALID_URL'
                    })
                    continue
                
                url = url.strip()
                logger.info(f"Processando URL {i+1}/{len(urls)}: {url}")
                
                # Baixar arquivo
                try:
                    downloaded_file = download_file_from_url(url, max_size)
                except Exception as e:
                    errors.append({
                        'index': i,
                        'url': url,
                        'error': str(e),
                        'error_code': 'DOWNLOAD_ERROR'
                    })
                    continue
                
                # Verificar tipo de arquivo
                filename = downloaded_file['filename']
                if not allowed_file(filename):
                    errors.append({
                        'index': i,
                        'url': url,
                        'filename': filename,
                        'error': f'Tipo de arquivo não suportado: {filename}',
                        'error_code': 'UNSUPPORTED_TYPE'
                    })
                    continue
                
                file_extension = filename.rsplit('.', 1)[1].lower()
                file_size = downloaded_file['size']
                
                # Validar tamanho
                if not validate_file_size(file_size, file_extension):
                    max_allowed = MAX_FILE_SIZES.get(file_extension, MAX_FILE_SIZES['default'])
                    errors.append({
                        'index': i,
                        'url': url,
                        'filename': filename,
                        'error': f'Arquivo muito grande. Máximo: {max_allowed // (1024*1024)}MB',
                        'error_code': 'FILE_TOO_LARGE'
                    })
                    continue
                
                # Processar arquivo
                with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
                    temp_file.write(downloaded_file['content'].getvalue())
                    temp_file_path = temp_file.name
                
                try:
                    # Extrair dados baseado na extensão
                    if file_extension == 'docx':
                        result = extract_text_from_docx(temp_file_path)
                    elif file_extension == 'pdf':
                        result = extract_text_from_pdf(temp_file_path)
                    elif file_extension in ['xlsx', 'xls']:
                        result = extract_data_from_excel(temp_file_path)
                    elif file_extension == 'csv':
                        result = extract_text_from_csv(temp_file_path)
                    elif file_extension == 'txt':
                        result = extract_text_from_txt(temp_file_path)
                    elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                        result = extract_text_from_image(temp_file_path)
                    else:
                        errors.append({
                            'index': i,
                            'url': url,
                            'filename': filename,
                            'error': f'Processamento para {file_extension} não implementado',
                            'error_code': 'PROCESSING_NOT_IMPLEMENTED'
                        })
                        continue
                    
                    processing_time = time.time() - start_time
                    total_processing_time += processing_time
                    
                    # Adicionar informações
                    result['file_info'] = {
                        'filename': filename,
                        'type': file_extension,
                        'mime_type': ALLOWED_EXTENSIONS.get(file_extension, 'unknown'),
                        'size_bytes': file_size,
                        'size_mb': round(file_size / (1024*1024), 2),
                        'processing_time': round(processing_time, 2)
                    }
                    
                    result['download_info'] = {
                        'source_url': url,
                        'content_type': downloaded_file['content_type'],
                        'download_successful': True,
                        'filename_from_url': filename
                    }
                    
                    results.append({
                        'index': i,
                        'url': url,
                        'filename': filename,
                        'success': True,
                        'data': result
                    })
                    
                    logger.info(f"URL processada com sucesso: {filename} de {url} ({file_size} bytes) em {processing_time:.2f}s")
                
                finally:
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                
            except Exception as e:
                processing_time = time.time() - start_time
                total_processing_time += processing_time
                
                logger.error(f"Erro ao processar URL {url}: {str(e)}")
                errors.append({
                    'index': i,
                    'url': url,
                    'error': str(e),
                    'error_code': 'PROCESSING_ERROR',
                    'processing_time': round(processing_time, 2)
                })
        
        # Estatísticas finais
        total_urls = len(urls)
        processed_count = len(results)
        failed_count = len(errors)
        success_rate = (processed_count / total_urls * 100) if total_urls > 0 else 0
        
        return jsonify({
            'success': True,
            'results': results,
            'errors': errors,
            'summary': {
                'total_urls': total_urls,
                'processed': processed_count,
                'failed': failed_count,
                'success_rate': round(success_rate, 1),
                'total_processing_time': round(total_processing_time, 2),
                'average_time_per_url': round(total_processing_time / total_urls, 2) if total_urls > 0 else 0
            }
        })
    
    except Exception as e:
        logger.error(f"Erro no processamento em lote de URLs: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e),
            'error_code': 'BULK_URL_PROCESSING_ERROR'
        }), 500 